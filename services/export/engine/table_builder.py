"""
Table builder — native PowerPoint tables.

Supports:
  - Alternating row background colors (alt row uses theme color)
  - Header row with theme header background + white text
  - Per-cell highlight from content.highlight_cells array
  - Change indicators: ▲ (up/green), ▼ (down/red), → (neutral/gray)
  - Column width distribution from headers[].width_percent
  - Cell emphasis (bold + highlight background)
  - Column alignment (left/center/right) from headers[].align

Schema fields consumed:
  content.table_title          → optional caption (not currently rendered — reserved)
  content.headers[].key        → column key for cell lookup
  content.headers[].label      → display label in header row
  content.headers[].width_percent → proportional column width
  content.headers[].align      → text alignment
  content.rows[].row_id        → stable row identifier (for highlight lookup)
  content.rows[].cells[key].value → display string
  content.rows[].cells[key].numeric_value → raw number (for formatting)
  content.rows[].cells[key].emphasis → bold + background
  content.rows[].cells[key].change_indicator → up/down/neutral/none
  content.highlight_cells      → [{row_id, column_key, reason}]
"""
from __future__ import annotations
from pptx.slide import Slide
from pptx.util import Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN

from engine.unit_converter import units_to_emu, font_units_to_pt, apply_font_scale
from engine.theme_resolver import ThemeTokens

_CHANGE_SYMBOLS = {
    "up":      "▲",
    "down":    "▼",
    "neutral": "→",
    "none":    "",
}

_CHANGE_COLORS = {
    "up":      (16, 185, 129),    # Emerald — positive
    "down":    (239, 68, 68),     # Red — negative
    "neutral": (148, 163, 184),   # Slate — flat
    "none":    None,              # Use body color
}


class TableBuilder:
    """
    Builds native PowerPoint table shapes.
    All text boxes inside the table are editable. The table can be resized
    and reformatted directly in PowerPoint.
    """

    def __init__(self, tokens: ThemeTokens, aspect_ratio: str = "16:9"):
        self.tokens       = tokens
        self.aspect_ratio = aspect_ratio

    def build_table(
        self,
        slide: Slide,
        headers: list[dict],
        rows:    list[dict],
        bounds:  dict,
        highlight_cells: list[dict],
        font_scale: float = 1.0,
    ) -> None:
        """
        Add a native PowerPoint table to the slide.

        Row count: len(rows) + 1 (header row).
        Column count: len(headers).
        """
        if not headers:
            return
        if not rows:
            return

        x = units_to_emu(bounds["x"],      'x',      self.aspect_ratio)
        y = units_to_emu(bounds["y"],      'y',      self.aspect_ratio)
        w = units_to_emu(bounds["width"],  'width',  self.aspect_ratio)
        h = units_to_emu(bounds["height"], 'height', self.aspect_ratio)
        w = max(w, 914_400)
        h = max(h, 457_200)

        num_cols = len(headers)
        num_rows = len(rows) + 1   # +1 for header

        graphic_frame = slide.shapes.add_table(num_rows, num_cols, x, y, w, h)
        table         = graphic_frame.table

        # ── Column widths from width_percent ──────────────────────────────────
        total_pct = sum(hdr.get("width_percent", 100.0 / num_cols) for hdr in headers)
        for ci, hdr in enumerate(headers):
            pct = hdr.get("width_percent", 100.0 / num_cols)
            table.columns[ci].width = int(w * (pct / total_pct))

        # ── Build highlight lookup set ────────────────────────────────────────
        highlight_set: set[tuple] = {
            (hc.get("row_id"), hc.get("column_key"))
            for hc in (highlight_cells or [])
            if hc.get("row_id") and hc.get("column_key")
        }

        # ── Header row (row index 0) ──────────────────────────────────────────
        header_font_pt = apply_font_scale(
            font_units_to_pt(18, self.aspect_ratio), font_scale
        )
        for ci, hdr in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = hdr.get("label", hdr.get("key", ""))

            cell.fill.solid()
            cell.fill.fore_color.rgb = RgbColor(*self.tokens.table_header_bg_rgb)

            p = cell.text_frame.paragraphs[0]
            p.font.bold      = True
            p.font.size      = Pt(header_font_pt)
            p.font.color.rgb = RgbColor(*self.tokens.table_header_text_rgb)
            p.font.name      = self.tokens.body_font_name
            p.alignment      = self._resolve_alignment(hdr.get("align", "left"))

        # ── Data rows (row index 1..n) ────────────────────────────────────────
        body_font_pt = apply_font_scale(
            font_units_to_pt(16, self.aspect_ratio), font_scale
        )
        for ri, row in enumerate(rows):
            row_id    = row.get("row_id", "")
            cells_map = row.get("cells", {})
            is_alt    = (ri % 2) == 1   # 0-indexed: row 0,2,4... → normal, 1,3,5... → alt

            for ci, hdr in enumerate(headers):
                key  = hdr.get("key", "")
                cell = table.cell(ri + 1, ci)

                cell_data    = cells_map.get(key, {})
                raw_value    = str(cell_data.get("value", ""))
                change_ind   = cell_data.get("change_indicator", "none")
                is_emphasized = bool(cell_data.get("emphasis", False))
                is_highlighted = (row_id, key) in highlight_set

                # Change indicator prefix
                symbol = _CHANGE_SYMBOLS.get(change_ind, "")
                display_text = f"{symbol} {raw_value}".strip() if symbol else raw_value

                cell.text = display_text

                # Cell background
                if is_highlighted:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RgbColor(*self.tokens.table_highlight_bg_rgb)
                elif is_alt:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RgbColor(*self.tokens.table_row_alt_bg_rgb)

                # Text formatting
                p = cell.text_frame.paragraphs[0]
                p.font.size      = Pt(body_font_pt)
                p.font.name      = self.tokens.body_font_name
                p.font.bold      = is_emphasized or is_highlighted
                p.alignment      = self._resolve_alignment(hdr.get("align", "left"))

                # Change indicator color overrides body color
                change_color = _CHANGE_COLORS.get(change_ind)
                p.font.color.rgb = RgbColor(*(change_color or self.tokens.body_rgb))

    @staticmethod
    def _resolve_alignment(align: str) -> PP_ALIGN:
        return {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(
            align, PP_ALIGN.LEFT
        )
