"""
Presentation renderer — main orchestrator.

Flow:
  1. Validation gate (raises ExportValidationError if deck has blocking errors)
  2. Resolve theme tokens
  3. Create python-pptx Presentation with correct slide dimensions
  4. For each slide (sorted by slide_index):
       a. Get LayoutSolution (warn + use defaults if missing)
       b. Route to correct builder class
       c. Call builder.build()
       d. Apply watermark if watermark_required
       e. Inject caveat footer if handling_caveats
       f. Add speaker notes to presenter view
       g. On ANY exception: log + add error placeholder slide (never crash entire export)
  5. Embed theme fonts
  6. Serialise Presentation to bytes
  7. Return bytes (caller uploads to S3)
"""
from __future__ import annotations
import logging
from io import BytesIO
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RgbColor

from engine.validation_gate import run_validation_gate
from engine.theme_resolver import resolve_theme
from engine.font_embedder import FontEmbedder
from engine.slide_builders.title_slide import TitleSlideBuilder
from engine.slide_builders.content_bullets import ContentBulletsBuilder
from engine.slide_builders.data_chart import DataChartBuilder
from engine.slide_builders.visual_split import VisualSplitBuilder
from engine.slide_builders.table_slide import TableSlideBuilder
from engine.slide_builders.section_divider import SectionDividerBuilder
from core.security import inject_watermark, inject_handling_caveats_footer
from core.exceptions import ExportValidationError, SlideExportError

logger = logging.getLogger(__name__)

# PowerPoint standard 16:9 slide dimensions in EMU
_SLIDE_W_EMU = 12_192_756
_SLIDE_H_EMU =  6_858_000

# Map slide_type → builder class
_BUILDER_MAP: dict[str, type] = {
    "title_slide":     TitleSlideBuilder,
    "content_bullets": ContentBulletsBuilder,
    "data_chart":      DataChartBuilder,
    "visual_split":    VisualSplitBuilder,
    "table":           TableSlideBuilder,
    "section_divider": SectionDividerBuilder,
}


class PresentationRenderer:
    """
    The only public interface of the export engine.
    Call render() with the full deck dict and layout_solutions dict.
    """

    def render(
        self,
        deck:             dict,
        layout_solutions: dict,
        plan_tier:        str = "free",
    ) -> bytes:
        """
        Build and return a .pptx file as raw bytes.

        Parameters
        ----------
        deck : complete deck JSON (from Phase 1, possibly with Phase 2 font_scale written back)
        layout_solutions : {slide_id: LayoutSolution dict} from Phase 2
        plan_tier : "free" | "pro" | "enterprise" (affects signed URL expiry, not rendering)

        Returns
        -------
        bytes — the complete .pptx binary. Upload to S3 immediately.

        Raises
        ------
        ExportValidationError
            If the deck has blocking_errors, missing required fields, or
            non-numeric chart values. The export must NOT proceed in these cases.
        """

        # ── Step 1: Validation gate ───────────────────────────────────────────
        is_valid, errors = run_validation_gate(deck)
        if not is_valid:
            raise ExportValidationError(errors)

        # ── Step 2: Resolve configuration ────────────────────────────────────
        theme_name    = deck.get("metadata", {}).get("theme", "modern_light")
        aspect_ratio  = deck.get("aspect_ratio", "16:9")
        tokens        = resolve_theme(theme_name)
        security      = deck.get("security_classification", {})
        watermark_req = security.get("watermark_required", False)
        class_level   = security.get("level", "internal")
        caveats       = security.get("handling_caveats", [])

        # ── Step 3: Create Presentation object ───────────────────────────────
        prs = Presentation()
        prs.slide_width  = Emu(_SLIDE_W_EMU)
        prs.slide_height = Emu(_SLIDE_H_EMU)

        # Always use blank slide layout (index 6) — we draw everything ourselves
        blank_layout = prs.slide_layouts[6]

        slides = sorted(
            deck.get("slides", []),
            key=lambda s: s.get("slide_index", 0),
        )
        total_slides     = len(slides)
        export_errors    = []   # Non-fatal slide errors collected here

        # ── Step 4: Build each slide ─────────────────────────────────────────
        for slide_data in slides:
            slide_id    = slide_data.get("slide_id", "unknown")
            slide_type  = slide_data.get("slide_type", "content_bullets")
            slide_index = slide_data.get("slide_index", 0)

            # Get LayoutSolution — warn if missing but continue
            layout_solution = layout_solutions.get(slide_id, {})
            if not layout_solution:
                logger.warning(
                    f"No LayoutSolution for slide {slide_id} (index {slide_index}). "
                    f"Builders will use element_bounds=None fallback paths."
                )

            try:
                pptx_slide = prs.slides.add_slide(blank_layout)

                # Route to correct builder
                BuilderClass = _BUILDER_MAP.get(slide_type)
                if not BuilderClass:
                    raise SlideExportError(
                        slide_id, slide_index,
                        f"Unknown slide_type '{slide_type}' — no builder registered."
                    )

                builder = BuilderClass(tokens, aspect_ratio)
                builder.build(
                    slide=pptx_slide,
                    slide_data=slide_data,
                    layout_solution=layout_solution,
                    slide_index=slide_index,
                    total_slides=total_slides,
                )

                # ── Security overlays ──────────────────────────────────────
                if watermark_req:
                    # Per-slide security can override deck-level classification
                    slide_level = (
                        slide_data.get("security_classification", {}).get("level")
                        or class_level
                    )
                    inject_watermark(pptx_slide, slide_level, caveats)
                elif caveats:
                    inject_handling_caveats_footer(pptx_slide, caveats)

                # ── Speaker notes ────────────────────────────────────────────
                speaker_notes = slide_data.get("speaker_notes", "").strip()
                if speaker_notes:
                    pptx_slide.notes_slide.notes_text_frame.text = speaker_notes

            except SlideExportError as exc:
                err_msg = str(exc)
                logger.error(f"SlideExportError: {err_msg}")
                export_errors.append(err_msg)
                self._add_error_placeholder(prs, blank_layout, slide_id, slide_index, err_msg)

            except Exception as exc:
                err_msg = f"Slide {slide_index} ({slide_id[:8]}): {exc}"
                logger.error(err_msg, exc_info=True)
                export_errors.append(err_msg)
                self._add_error_placeholder(prs, blank_layout, slide_id, slide_index, str(exc))

        # ── Step 5: Font embedding ───────────────────────────────────────────
        FontEmbedder().embed_fonts(prs, tokens)

        # ── Step 6: Serialise ────────────────────────────────────────────────
        output = BytesIO()
        prs.save(output)
        output.seek(0)

        pptx_bytes = output.getvalue()

        if len(pptx_bytes) > (100 * 1024 * 1024):  # 100 MB
            logger.warning(
                f"PPTX file size {len(pptx_bytes) / 1e6:.1f} MB exceeds 100 MB limit. "
                f"Consider reducing image assets or slide count."
            )

        if export_errors:
            logger.warning(
                f"Export completed with {len(export_errors)} slide error(s). "
                f"Error placeholder slides have been inserted."
            )

        return pptx_bytes

    def _add_error_placeholder(
        self,
        prs,
        layout,
        slide_id:    str,
        slide_index: int,
        error_msg:   str,
    ) -> None:
        """
        Add a visually distinct error placeholder slide at the expected position.

        This ensures the slide count is preserved. The user sees a slide labelled
        "Slide N failed to export" with the error message — they know what happened
        and which slide needs attention.

        NEVER silently skip a failed slide — always insert the placeholder.
        """
        slide = prs.slides.add_slide(layout)

        # Light red background
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = RgbColor(255, 240, 240)

        # Error title
        txBox = slide.shapes.add_textbox(
            Emu(500_000), Emu(500_000),
            Emu(11_192_756), Emu(2_000_000),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.clear()

        p = tf.paragraphs[0]
        p.text = f"Slide {slide_index + 1} — Export Failed"
        p.font.size      = Pt(24)
        p.font.bold      = True
        p.font.color.rgb = RgbColor(239, 68, 68)
        p.font.name      = "Inter"

        # Error message
        txBox2 = slide.shapes.add_textbox(
            Emu(500_000), Emu(2_700_000),
            Emu(11_192_756), Emu(2_000_000),
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        tf2.clear()

        p2 = tf2.paragraphs[0]
        p2.text = f"Error: {error_msg[:400]}"
        p2.font.size      = Pt(14)
        p2.font.color.rgb = RgbColor(127, 29, 29)
        p2.font.name      = "Inter"

        # Slide ID for debugging
        txBox3 = slide.shapes.add_textbox(
            Emu(500_000), Emu(4_900_000),
            Emu(11_192_756), Emu(600_000),
        )
        tf3 = txBox3.text_frame
        tf3.clear()
        p3 = tf3.paragraphs[0]
        p3.text = f"slide_id: {slide_id}"
        p3.font.size      = Pt(10)
        p3.font.color.rgb = RgbColor(180, 100, 100)
        p3.font.name      = "Inter"
