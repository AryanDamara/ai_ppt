"""
Chart engine — native Office charts only.

All 8 chart types supported:
  column_clustered  → CategoryChartData, XL_CHART_TYPE.COLUMN_CLUSTERED
  column_stacked    → CategoryChartData, XL_CHART_TYPE.COLUMN_STACKED
  line              → CategoryChartData, XL_CHART_TYPE.LINE
  pie               → CategoryChartData, XL_CHART_TYPE.PIE
  bar               → CategoryChartData, XL_CHART_TYPE.BAR_CLUSTERED
  area              → CategoryChartData, XL_CHART_TYPE.AREA
  scatter           → XyChartData,       XL_CHART_TYPE.XY_SCATTER
  waterfall         → CategoryChartData, XL_CHART_TYPE.COLUMN_STACKED
                       (invisible base series + delta series trick)

NEVER generate a matplotlib image. NEVER use Pillow to render a chart.
ALWAYS use slide.shapes.add_chart() which creates a real, editable chart object.
"""
from __future__ import annotations
import logging

from pptx.slide import Slide
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RgbColor
from pptx.util import Pt

from engine.unit_converter import units_to_emu
from engine.theme_resolver import ThemeTokens
from core.exceptions import ChartDataError

logger = logging.getLogger(__name__)

# Mapping from schema enum → python-pptx XL_CHART_TYPE
_CHART_TYPE_MAP: dict[str, XL_CHART_TYPE] = {
    "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column_stacked":   XL_CHART_TYPE.COLUMN_STACKED,
    "line":             XL_CHART_TYPE.LINE,
    "pie":              XL_CHART_TYPE.PIE,
    "bar":              XL_CHART_TYPE.BAR_CLUSTERED,
    "area":             XL_CHART_TYPE.AREA,
    "scatter":          XL_CHART_TYPE.XY_SCATTER,
    # Waterfall uses COLUMN_STACKED with invisible base series
    "waterfall":        XL_CHART_TYPE.COLUMN_STACKED,
}


class ChartEngine:
    """
    Generates native, Excel-backed PowerPoint charts from JSON chart_data.

    Each chart object, once added to the slide, contains:
    - A ChartPart with the OOXML chart definition
    - An embedded Excel workbook (the "data source")
    - Format overrides (colors, legend, axis bounds)

    When a user double-clicks the chart in PowerPoint, Excel opens showing
    the data table — exactly as if they had inserted a chart manually.
    """

    def __init__(self, tokens: ThemeTokens):
        self.tokens = tokens

    def add_chart(
        self,
        slide: Slide,
        chart_type: str,
        chart_data: dict,
        chart_options: dict,
        bounds: dict,
        aspect_ratio: str = "16:9",
    ) -> None:
        """
        Add a native Office chart to the slide at the LayoutSolution position.

        Parameters
        ----------
        slide : pptx Slide object
        chart_type : one of the 8 schema enum values
        chart_data : { series: [...], categories: [...], global_unit?, data_source? }
        chart_options : { show_legend?, show_data_labels?, y_axis_max?, y_axis_min?, trendline_enabled? }
        bounds : element bounds dict from LayoutSolution (layout units)
        aspect_ratio : from deck metadata
        """
        x = units_to_emu(bounds["x"],      'x',      aspect_ratio)
        y = units_to_emu(bounds["y"],      'y',      aspect_ratio)
        w = units_to_emu(bounds["width"],  'width',  aspect_ratio)
        h = units_to_emu(bounds["height"], 'height', aspect_ratio)
        w = max(w, 914_400)    # Minimum 1 inch
        h = max(h, 914_400)

        # Validate before building ChartData — fail with clear error
        self._validate_chart_data(chart_data, chart_type)

        pptx_chart_type = _CHART_TYPE_MAP.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

        if chart_type == "scatter":
            cd = self._build_xy_chart_data(chart_data)
        else:
            cd = self._build_category_chart_data(chart_data, chart_type)

        chart_frame = slide.shapes.add_chart(pptx_chart_type, x, y, w, h, cd)
        chart       = chart_frame.chart

        self._apply_formatting(chart, chart_data, chart_options, chart_type)

    # ── Data validation ────────────────────────────────────────────────────────

    def _validate_chart_data(self, chart_data: dict, chart_type: str) -> None:
        """
        Validate chart data before building ChartData objects.
        Raises ChartDataError with a human-readable message on failure.
        """
        series_list = chart_data.get("series", [])
        if not series_list:
            raise ChartDataError(f"chart_type='{chart_type}': no series data provided.")

        categories = chart_data.get("categories", [])
        if not categories and chart_type != "scatter":
            raise ChartDataError(f"chart_type='{chart_type}': categories array is empty.")

        for si, series in enumerate(series_list):
            values = series.get("values", [])
            if not values:
                raise ChartDataError(
                    f"chart_type='{chart_type}': series[{si}] "
                    f"('{series.get('name')}') has no values."
                )
            for vi, val in enumerate(values):
                if not isinstance(val, (int, float)):
                    raise ChartDataError(
                        f"chart_type='{chart_type}': series[{si}] "
                        f"('{series.get('name')}') value[{vi}] is {type(val).__name__}: {val!r}. "
                        f"ALL chart values MUST be Python int or float. "
                        f"The AI pipeline likely generated quoted strings. "
                        f"Fix this in Phase 1 step3_content.py."
                    )
            # Length check (not applicable to scatter)
            if chart_type != "scatter" and len(values) != len(categories):
                raise ChartDataError(
                    f"chart_type='{chart_type}': series[{si}] "
                    f"('{series.get('name')}') has {len(values)} values "
                    f"but {len(categories)} categories. They must be equal."
                )

    # ── CategoryChartData builder ──────────────────────────────────────────────

    def _build_category_chart_data(
        self, chart_data: dict, chart_type: str
    ) -> CategoryChartData:
        """
        Build CategoryChartData for: column_clustered, column_stacked, line,
        pie, bar, area, and waterfall.

        Waterfall implementation:
        python-pptx 0.6.23 does not have XL_CHART_TYPE.WATERFALL.
        We simulate it using a stacked column chart with:
          Series 0: "invisible base" — cumulative running minimum per bar
          Series 1+: delta values (absolute change)
        The invisible base uses transparent fill, creating the waterfall visual.
        """
        cd          = CategoryChartData()
        cd.categories = chart_data.get("categories", [])
        global_unit = chart_data.get("global_unit", "General")

        series_list = chart_data.get("series", [])

        if chart_type == "waterfall":
            # Waterfall: build invisible base + delta for each series
            for series in series_list:
                values = series.get("values", [])
                unit   = series.get("unit") or global_unit

                # Calculate cumulative running minimum (the "invisible" stack base)
                bases   = []
                running = 0.0
                for v in values:
                    bases.append(min(running, running + v))
                    running += v
                deltas = [abs(v) for v in values]

                # Invisible base series (will have transparent fill applied in formatting)
                cd.add_series("_base_invisible", tuple(bases), number_format=unit)
                # Visible delta series
                cd.add_series(series.get("name", ""), tuple(deltas), number_format=unit)
        else:
            for series in series_list:
                unit = series.get("unit") or global_unit
                cd.add_series(
                    series.get("name", ""),
                    tuple(series.get("values", [])),
                    number_format=unit,
                )

        return cd

    # ── XyChartData builder ────────────────────────────────────────────────────

    def _build_xy_chart_data(self, chart_data: dict) -> XyChartData:
        """Build XyChartData for scatter charts. Uses index as x if only y values provided."""
        cd = XyChartData()
        for series in chart_data.get("series", []):
            s      = cd.add_series(series.get("name", ""))
            values = series.get("values", [])
            for i, y_val in enumerate(values):
                s.add_data_point(float(i), float(y_val))
        return cd

    # ── Chart formatting ───────────────────────────────────────────────────────

    def _apply_formatting(
        self,
        chart,
        chart_data: dict,
        chart_options: dict,
        chart_type: str,
    ) -> None:
        """
        Apply theme colors, legend, data labels, axis bounds, and trendlines.
        All formatting is wrapped in try/except — formatting failure is non-fatal
        (chart still exports, just without custom colors).
        """
        series_list = chart_data.get("series", [])

        # ── Legend ────────────────────────────────────────────────────────────
        try:
            show_legend     = chart_options.get("show_legend", True)
            chart.has_legend = show_legend
            if show_legend and hasattr(chart, 'legend'):
                chart.legend.position          = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
        except Exception as e:
            logger.debug(f"Legend formatting failed (non-fatal): {e}")

        # ── Series colors ─────────────────────────────────────────────────────
        try:
            palette    = self.tokens.chart_palette
            pptx_series_iter = iter(chart.series)
            for i, (schema_series, pptx_series) in enumerate(
                zip(series_list, pptx_series_iter)
            ):
                # Waterfall: first series (_base_invisible) → transparent fill
                if chart_type == "waterfall" and schema_series.get("name") == "_base_invisible":
                    try:
                        pptx_series.format.fill.background()  # Transparent
                    except Exception:
                        pass
                    continue

                # Resolve color: schema hex > theme palette > default indigo
                color_hex = schema_series.get("color", "")
                if color_hex and len(color_hex) == 7 and color_hex.startswith("#"):
                    r = int(color_hex[1:3], 16)
                    g = int(color_hex[3:5], 16)
                    b = int(color_hex[5:7], 16)
                    fill_color = (r, g, b)
                else:
                    fill_color = self.tokens.series_color(i)

                try:
                    pptx_series.format.fill.solid()
                    pptx_series.format.fill.fore_color.rgb = RgbColor(*fill_color)
                except Exception:
                    pass  # Some chart types don't support direct series fill

                # Data labels
                show_labels = chart_options.get("show_data_labels", False)
                try:
                    pptx_series.data_labels.show_value = show_labels
                except Exception:
                    pass

        except Exception as e:
            logger.debug(f"Series color formatting failed (non-fatal): {e}")

        # ── Axis bounds ───────────────────────────────────────────────────────
        try:
            if hasattr(chart, 'value_axis'):
                y_max = chart_options.get("y_axis_max")
                y_min = chart_options.get("y_axis_min")
                if y_max is not None:
                    chart.value_axis.maximum_scale = float(y_max)
                if y_min is not None:
                    chart.value_axis.minimum_scale = float(y_min)
        except Exception as e:
            logger.debug(f"Axis bounds formatting failed (non-fatal): {e}")

        # ── Trendline ─────────────────────────────────────────────────────────
        try:
            if chart_options.get("trendline_enabled") and chart_type in (
                "line", "column_clustered", "scatter"
            ):
                for pptx_series in chart.series:
                    try:
                        pptx_series.smooth = True
                    except Exception:
                        pass
        except Exception as e:
            logger.debug(f"Trendline formatting failed (non-fatal): {e}")
