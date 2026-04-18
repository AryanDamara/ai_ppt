"""
visual_split builder.

Schema fields consumed:
  action_title                    → rendered as slide title in the text zone
  content.supporting_text         → body text (2-4 sentences)
  content.image_asset_id          → UUID matching assets[].asset_id
  content.image_keyword           → fallback placeholder label if asset missing
  content.text_position           → "left" | "right" | "overlay"
  content.image_treatment         → "original"|"monochrome"|"duotone"|"gradient_overlay"
  assets[].asset_id               → matched to content.image_asset_id
  assets[].source_uri             → S3 URL or presigned URL for image bytes
  assets[].variants.web_optimized → preferred image variant

LayoutSolution zones consumed:
  "title"   → action_title text box (in the text half)
  "body"    → supporting_text text box
  "image"   → image bounding box (in the image half)
"""
from __future__ import annotations
from pptx.slide import Slide
from pptx.util import Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO
from .base_builder import BaseSlideBuilder
from engine.unit_converter import units_to_emu


class VisualSplitBuilder(BaseSlideBuilder):

    def build(
        self,
        slide: Slide,
        slide_data: dict,
        layout_solution: dict,
        slide_index: int,
        total_slides: int,
    ) -> None:
        self.set_slide_background(slide)

        content    = slide_data.get("content", {})
        assets     = slide_data.get("assets", [])
        font_scale = self.get_font_scale(slide_data)
        overrides  = slide_data.get("template", {}).get("overrides", {})
        is_rtl     = self.is_rtl_layout(layout_solution)
        align      = self.get_pptx_alignment("left", is_rtl)

        # ── Image zone ────────────────────────────────────────────────────────
        image_bounds = self.get_element_bounds(layout_solution, "image")
        if image_bounds:
            self._render_image(slide, content, assets, image_bounds)

        # ── Title ─────────────────────────────────────────────────────────────
        title_bounds = self.get_element_bounds(layout_solution, "title")
        if title_bounds:
            self.add_textbox(
                slide, title_bounds,
                slide_data.get("action_title", ""),
                font_size_units=title_bounds.get("font_size_units", 48),
                font_scale=font_scale,
                color_rgb=self.tokens.title_rgb,
                bold=True,
                font_name=self.tokens.display_font_name,
                alignment=align,
            )

        # ── Supporting text ───────────────────────────────────────────────────
        body_bounds     = self.get_element_bounds(layout_solution, "body")
        supporting_text = content.get("supporting_text", "")
        if supporting_text and body_bounds:
            self.add_textbox(
                slide, body_bounds, supporting_text,
                font_size_units=body_bounds.get("font_size_units", 24),
                font_scale=font_scale,
                color_rgb=self.tokens.body_rgb,
                font_name=self.tokens.body_font_name,
                alignment=align,
                word_wrap=True,
            )

        self.add_footer(
            slide,
            source_footer=slide_data.get("source_footer"),
            slide_index=slide_index,
            total_slides=total_slides,
            hide_footer=overrides.get("hide_footer", False),
            hide_slide_number=overrides.get("hide_slide_number", False),
        )

    def _render_image(
        self,
        slide: Slide,
        content: dict,
        assets: list,
        image_bounds: dict,
    ) -> None:
        """
        Attempt to fetch and embed the image from S3.
        If fetch fails for any reason, render a styled placeholder rectangle
        showing the image_keyword as label text.

        The placeholder is NOT a silent failure — it is visible in the PPTX
        so the user knows an image was intended here.
        """
        x = units_to_emu(image_bounds["x"],     'x',      self.aspect_ratio)
        y = units_to_emu(image_bounds["y"],     'y',      self.aspect_ratio)
        w = units_to_emu(image_bounds["width"], 'width',  self.aspect_ratio)
        h = units_to_emu(image_bounds["height"],'height', self.aspect_ratio)
        w = max(w, 91_440)
        h = max(h, 91_440)

        image_asset_id = content.get("image_asset_id")
        image_bytes    = None

        if image_asset_id:
            asset = next(
                (a for a in assets if a.get("asset_id") == image_asset_id), None
            )
            if asset:
                try:
                    from engine.image_processor import ImageProcessor
                    src = (
                        asset.get("source_uri") or
                        asset.get("variants", {}).get("web_optimized") or
                        asset.get("variants", {}).get("print_300dpi")
                    )
                    if src:
                        proc       = ImageProcessor()
                        raw        = proc.fetch(src)
                        treatment  = content.get("image_treatment", "original")
                        image_bytes = proc.apply_treatment(raw, treatment)
                except Exception as exc:
                    import logging
                    logging.getLogger(__name__).warning(
                        f"Image fetch failed for asset {image_asset_id}: {exc}"
                    )

        if image_bytes:
            slide.shapes.add_picture(BytesIO(image_bytes), x, y, w, h)
        else:
            # Placeholder rectangle with keyword label
            shape = slide.shapes.add_shape(1, x, y, w, h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RgbColor(*self.tokens.table_header_bg_rgb)
            shape.line.fill.background()

            keyword = content.get("image_keyword", "Image")
            tf = shape.text_frame
            tf.word_wrap = True
            tf.clear()
            p   = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"[ {keyword} ]"
            run.font.size      = Pt(14)
            run.font.italic    = True
            run.font.color.rgb = RgbColor(*self.tokens.footer_rgb)
            run.font.name      = self.tokens.body_font_name
