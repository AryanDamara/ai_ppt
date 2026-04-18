"""
Base slide builder.

All 6 slide type builders inherit from this class.
Provides shared helpers so builder files stay focused on slide-type logic.

The single most important rule: EVERY position, size, and font size comes from
the LayoutSolution via units_to_emu(). Never hardcode a number in pixels, inches,
or EMU in any builder subclass.
"""
from __future__ import annotations
from pptx.slide import Slide
from pptx.util import Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.shapes.base import BaseShape
from typing import Optional

from engine.unit_converter import (
    units_to_emu, font_units_to_pt, apply_font_scale,
    FOOTER_ZONE_Y_UNITS, SAFE_AREA_MARGIN_LEFT, SAFE_AREA_MARGIN_RIGHT, SLIDE_WIDTH_UNITS,
)
from engine.theme_resolver import ThemeTokens


class BaseSlideBuilder:
    """
    Shared utilities for all slide type builders.
    Do not instantiate directly — subclass for each slide_type.
    """

    def __init__(self, tokens: ThemeTokens, aspect_ratio: str = "16:9"):
        self.tokens       = tokens
        self.aspect_ratio = aspect_ratio

    # ── Core helper: add a text box ───────────────────────────────────────────

    def add_textbox(
        self,
        slide: Slide,
        element_bounds: dict,
        text: str,
        font_size_units: int,
        font_scale:    float          = 1.0,
        color_rgb:     tuple          = (0, 0, 0),
        bold:          bool           = False,
        italic:        bool           = False,
        underline:     bool           = False,
        font_name:     str            = "Inter",
        alignment:     PP_ALIGN       = PP_ALIGN.LEFT,
        word_wrap:     bool           = True,
        vertical_anchor: MSO_ANCHOR   = MSO_ANCHOR.TOP,
    ) -> BaseShape:
        """
        Add a single-paragraph text box positioned from LayoutSolution bounds.

        Parameters
        ----------
        slide : pptx Slide object
        element_bounds : dict with keys x, y, width, height (layout units)
        text : display string. If empty string, still adds the shape (blank).
        font_size_units : from LayoutSolution element['font_size_units']
        font_scale : from slide.template.overrides.font_scale (default 1.0)
        color_rgb : (r, g, b) tuple
        bold, italic, underline : text run styling
        font_name : must match an embedded or system font name
        alignment : PP_ALIGN.LEFT / CENTER / RIGHT
        word_wrap : True for body text, False for slide numbers
        vertical_anchor : MSO_ANCHOR.TOP / MIDDLE / BOTTOM

        Returns
        -------
        The created shape object (rarely needed by callers).

        Note
        ----
        The font_scale is applied via apply_font_scale() which clamps to [0.8, 1.2]
        before multiplying. Never call apply_font_scale() on the same value twice.
        """
        x = units_to_emu(element_bounds["x"],     'x',      self.aspect_ratio)
        y = units_to_emu(element_bounds["y"],     'y',      self.aspect_ratio)
        w = units_to_emu(element_bounds["width"], 'width',  self.aspect_ratio)
        h = units_to_emu(element_bounds["height"],'height', self.aspect_ratio)

        # Guard: python-pptx throws on zero-dimension shapes
        w = max(w, 91_440)   # Minimum 0.1 inch in EMU
        h = max(h, 91_440)

        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf    = txBox.text_frame
        tf.word_wrap       = word_wrap
        tf.vertical_anchor = vertical_anchor
        tf.clear()

        p   = tf.paragraphs[0]
        p.alignment = alignment

        run = p.add_run()
        run.text = text or ""

        font_pt = font_units_to_pt(font_size_units, self.aspect_ratio)
        font_pt = apply_font_scale(font_pt, font_scale)

        run.font.size       = Pt(font_pt)
        run.font.bold       = bold
        run.font.italic     = italic
        run.font.underline  = underline
        run.font.name       = font_name
        run.font.color.rgb  = RgbColor(*color_rgb)

        return txBox

    # ── Core helper: multi-paragraph text box ────────────────────────────────

    def add_textbox_multiline(
        self,
        slide: Slide,
        element_bounds: dict,
        paragraphs: list[dict],
        font_scale: float       = 1.0,
        alignment:  PP_ALIGN    = PP_ALIGN.LEFT,
    ) -> BaseShape:
        """
        Add a text box with multiple styled paragraphs.

        Each entry in paragraphs is a dict with:
          text           : str (required)
          font_size_units: int (required, from LayoutSolution or explicit)
          bold           : bool (default False)
          italic         : bool (default False)
          color_rgb      : tuple (default self.tokens.body_rgb)
          font_name      : str (default self.tokens.body_font_name)
          indent_level   : int 0/1/2 (default 0)
          space_after_pt : float points of space after this paragraph (default 0)

        font_scale is applied to ALL paragraphs — it is the slide-level scale factor
        from template.overrides.font_scale.
        """
        x = units_to_emu(element_bounds["x"],     'x',      self.aspect_ratio)
        y = units_to_emu(element_bounds["y"],     'y',      self.aspect_ratio)
        w = units_to_emu(element_bounds["width"], 'width',  self.aspect_ratio)
        h = units_to_emu(element_bounds["height"],'height', self.aspect_ratio)
        w = max(w, 91_440)
        h = max(h, 91_440)

        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf    = txBox.text_frame
        tf.word_wrap = True
        tf.clear()

        for idx, spec in enumerate(paragraphs):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.alignment = alignment

            indent = spec.get("indent_level", 0)
            if indent > 0:
                # Level-based indent: python-pptx pPr 'lvl' attribute
                pPr = p._p.get_or_add_pPr()
                pPr.set("lvl", str(indent))

            space_after = spec.get("space_after_pt", 0.0)
            if space_after > 0:
                from pptx.util import Pt as _Pt
                p.space_after = _Pt(space_after)

            run = p.add_run()
            run.text = spec.get("text", "")

            font_pt = font_units_to_pt(spec.get("font_size_units", 24), self.aspect_ratio)
            font_pt = apply_font_scale(font_pt, font_scale)

            run.font.size      = Pt(font_pt)
            run.font.bold      = spec.get("bold", False)
            run.font.italic    = spec.get("italic", False)
            run.font.name      = spec.get("font_name", self.tokens.body_font_name)
            run.font.color.rgb = RgbColor(*spec.get("color_rgb", self.tokens.body_rgb))

        return txBox

    # ── Slide background ──────────────────────────────────────────────────────

    def set_slide_background(self, slide: Slide) -> None:
        """
        Set the slide background to the theme's background color.
        Called at the START of every builder's .build() method.
        """
        bg   = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RgbColor(*self.tokens.slide_background_rgb)

    # ── Footer zone ───────────────────────────────────────────────────────────

    def add_footer(
        self,
        slide: Slide,
        source_footer:      Optional[str],
        slide_index:        int,
        total_slides:       int,
        hide_footer:        bool = False,
        hide_slide_number:  bool = False,
    ) -> None:
        """
        Render the footer zone: source text on the left, slide number on the right.

        Footer zone: y = FOOTER_ZONE_Y_UNITS (710), height = remainder of slide.
        Both text boxes use font_size_units=14 (caption-size text).

        hide_footer suppresses the entire footer zone (title_slide default).
        hide_slide_number suppresses only the slide number.
        """
        if hide_footer:
            return

        slide_h_units = {"16:9": 562, "4:3": 750}[self.aspect_ratio]
        footer_h_units = max(10, slide_h_units - FOOTER_ZONE_Y_UNITS)

        if source_footer:
            self.add_textbox(
                slide,
                element_bounds={"x": SAFE_AREA_MARGIN_LEFT, "y": FOOTER_ZONE_Y_UNITS,
                                 "width": 700, "height": footer_h_units},
                text=source_footer,
                font_size_units=14,
                color_rgb=self.tokens.footer_rgb,
                font_name=self.tokens.body_font_name,
                alignment=PP_ALIGN.LEFT,
                word_wrap=False,
            )

        if not hide_slide_number:
            self.add_textbox(
                slide,
                element_bounds={"x": 820, "y": FOOTER_ZONE_Y_UNITS,
                                 "width": 120, "height": footer_h_units},
                text=f"{slide_index + 1} / {total_slides}",
                font_size_units=14,
                color_rgb=self.tokens.footer_rgb,
                font_name=self.tokens.body_font_name,
                alignment=PP_ALIGN.RIGHT,
                word_wrap=False,
            )

    # ── Utility helpers ───────────────────────────────────────────────────────

    def get_font_scale(self, slide_data: dict) -> float:
        """
        Read font_scale from template.overrides.font_scale.
        Written by Phase 2 tier-2 constraint relaxation when content overflowed.
        Always clamped to [0.8, 1.2]. Default is 1.0 (no scaling).
        """
        raw = slide_data.get("template", {}).get("overrides", {}).get("font_scale", 1.0)
        try:
            return max(0.8, min(1.2, float(raw)))
        except (TypeError, ValueError):
            return 1.0

    def get_element_bounds(
        self,
        layout_solution: dict,
        element_id: str,
    ) -> Optional[dict]:
        """
        Safely retrieve element bounds from a LayoutSolution.

        Returns None if element_id is not present in the solution.
        Callers MUST check for None and skip rendering if missing.
        This is non-fatal — a missing element means the layout engine
        decided not to place that element (e.g. empty optional zone).

        Usage pattern:
            bounds = self.get_element_bounds(layout_solution, "title")
            if bounds:
                self.add_textbox(slide, bounds, ...)
        """
        return layout_solution.get("elements", {}).get(element_id)

    def get_pptx_alignment(self, text_align: str, is_rtl: bool = False) -> PP_ALIGN:
        """
        Convert layout engine text_align string to python-pptx PP_ALIGN constant.
        RTL text defaults to right alignment regardless of text_align value.
        """
        if is_rtl:
            return PP_ALIGN.RIGHT
        return {
            "left":   PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right":  PP_ALIGN.RIGHT,
        }.get(text_align, PP_ALIGN.LEFT)

    def is_rtl_layout(self, layout_solution: dict) -> bool:
        """True if Phase 2 determined this slide uses right-to-left layout."""
        return bool(layout_solution.get("is_rtl", False))
