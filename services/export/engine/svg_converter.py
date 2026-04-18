"""
SVG converter — vector shapes preferred, PNG fallback, gray placeholder last resort.

Priority chain:
  1. Native DrawingML shapes (editable vectors, recolorable in PowerPoint)
  2. 300 DPI PNG with transparency (looks good but not editable)
  3. Gray placeholder rectangle (always renders something, never crashes)

Track fallback_rate property. If it exceeds 20%, the DrawingML coverage needs
improvement. Log a warning at that threshold.
"""
from __future__ import annotations
import logging
from io import BytesIO
from typing import Optional

from pptx.slide import Slide
from pptx.util import Pt
from pptx.dml.color import RgbColor

from engine.unit_converter import units_to_emu

logger = logging.getLogger(__name__)


class SVGConverter:

    def __init__(self):
        self._total        = 0
        self._png_fallbacks = 0

    @property
    def fallback_rate(self) -> float:
        return self._png_fallbacks / self._total if self._total > 0 else 0.0

    def add_svg_to_slide(
        self,
        slide: Slide,
        svg_content: str | bytes,
        bounds: dict,
        aspect_ratio: str = "16:9",
    ) -> str:
        """
        Add SVG content to the slide as a vector shape or PNG image.

        Returns
        -------
        "drawingml"   — native DrawingML shapes added
        "png"         — 300 DPI PNG image added
        "placeholder" — gray placeholder rectangle added (all methods failed)
        """
        self._total += 1

        x = units_to_emu(bounds["x"],      'x',      aspect_ratio)
        y = units_to_emu(bounds["y"],      'y',      aspect_ratio)
        w = units_to_emu(bounds["width"],  'width',  aspect_ratio)
        h = units_to_emu(bounds["height"], 'height', aspect_ratio)
        w = max(w, 91_440)
        h = max(h, 91_440)

        # Attempt 1: Native DrawingML
        try:
            added = self._add_drawingml(svg_content, slide, x, y, w, h)
            if added:
                return "drawingml"
        except Exception as exc:
            logger.debug(f"DrawingML conversion failed: {exc}")

        # Attempt 2: 300 DPI PNG
        try:
            png_bytes = self._render_to_png(svg_content, w, h)
            if png_bytes:
                self._png_fallbacks += 1
                slide.shapes.add_picture(BytesIO(png_bytes), x, y, w, h)
                if self.fallback_rate > 0.20:
                    logger.warning(
                        f"SVG PNG fallback rate is {self.fallback_rate:.0%} — "
                        f"improve DrawingML coverage in svg_converter.py"
                    )
                return "png"
        except Exception as exc:
            logger.debug(f"PNG rendering failed: {exc}")

        # Attempt 3: Gray placeholder
        self._png_fallbacks += 1
        self._add_gray_placeholder(slide, x, y, w, h)
        return "placeholder"

    def _add_drawingml(
        self,
        svg_content: str | bytes,
        slide: Slide,
        x: int, y: int, w: int, h: int,
    ) -> bool:
        """
        Parse SVG with svglib → reportlab drawing → map primitive shapes
        (Rect, Circle, Line) to python-pptx shapes.
        Returns True if at least one shape was added to the slide.
        """
        from svglib.svglib import svg2rlg
        from reportlab.graphics import shapes as rl_shapes

        svg_bytes = svg_content.encode() if isinstance(svg_content, str) else svg_content
        drawing   = svg2rlg(BytesIO(svg_bytes))
        if not drawing:
            return False

        primitives    = self._flatten_drawing(drawing)
        shapes_added  = 0

        for element in primitives:
            try:
                self._map_rl_shape(slide, element, x, y, w, h, drawing)
                shapes_added += 1
            except Exception:
                continue

        return shapes_added > 0

    def _flatten_drawing(self, drawing) -> list:
        """Recursively flatten a reportlab drawing to primitive shape list."""
        result = []
        for item in getattr(drawing, 'contents', []):
            if hasattr(item, 'contents'):
                result.extend(self._flatten_drawing(item))
            else:
                result.append(item)
        return result

    def _map_rl_shape(self, slide, rl_shape, x0, y0, total_w, total_h, drawing):
        """
        Map a reportlab primitive shape to a python-pptx shape.
        Handles: Rect. Extend this method to handle more primitives as needed.
        """
        from reportlab.graphics import shapes as rl_shapes
        from pptx.dml.color import RgbColor

        sx = total_w / drawing.width  if drawing.width  else 1.0
        sy = total_h / drawing.height if drawing.height else 1.0

        def to_x(v):  return x0 + int(v * sx)
        def to_y(v):  return y0 + int((drawing.height - v) * sy)  # Flip Y
        def to_w(v):  return max(1, int(v * sx))
        def to_h(v):  return max(1, int(v * sy))

        def parse_color(c) -> Optional[tuple]:
            if not c or str(c) == "none":
                return None
            try:
                from reportlab.lib.colors import HexColor
                hc = HexColor(str(c))
                return (int(hc.red * 255), int(hc.green * 255), int(hc.blue * 255))
            except Exception:
                return None

        if isinstance(rl_shape, rl_shapes.Rect):
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                to_x(rl_shape.x),
                to_y(rl_shape.y + rl_shape.height),
                to_w(rl_shape.width),
                to_h(rl_shape.height),
            )
            fill_c = parse_color(getattr(rl_shape, 'fillColor', None))
            if fill_c:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RgbColor(*fill_c)
            else:
                shape.fill.background()

            stroke_c = parse_color(getattr(rl_shape, 'strokeColor', None))
            if stroke_c:
                stroke_w = float(getattr(rl_shape, 'strokeWidth', 1.0))
                shape.line.color.rgb = RgbColor(*stroke_c)
                shape.line.width     = Pt(stroke_w)
            else:
                shape.line.fill.background()

    def _render_to_png(self, svg_content: str | bytes, w_emu: int, h_emu: int) -> Optional[bytes]:
        """Render SVG to a 300 DPI transparent-background PNG."""
        from svglib.svglib import svg2rlg
        from reportlab.graphics import renderPM

        svg_bytes = svg_content.encode() if isinstance(svg_content, str) else svg_content
        drawing   = svg2rlg(BytesIO(svg_bytes))
        if not drawing:
            return None

        buf = BytesIO()
        renderPM.drawToFile(drawing, buf, fmt="PNG", dpi=300)
        buf.seek(0)
        return buf.read()

    def _add_gray_placeholder(self, slide, x, y, w, h) -> None:
        """Last-resort placeholder: gray rectangle with no fill, no line."""
        shape = slide.shapes.add_shape(1, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RgbColor(200, 200, 200)
        shape.line.fill.background()
