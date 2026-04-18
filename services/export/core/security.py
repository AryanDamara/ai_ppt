"""
Security: watermark injection and handling caveat footer.

Watermarks are diagonal text boxes placed at z-order above all content.
The rotation attribute on the TextBox shape achieves the diagonal effect.
"""
from pptx.slide import Slide
from pptx.util import Pt, Emu
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN

# Standard slide dimensions in EMU (16:9 PowerPoint default)
_SLIDE_W_EMU = 12_192_756
_SLIDE_H_EMU =  6_858_000

_WATERMARK_LABELS: dict[str, str] = {
    "public":       "PUBLIC",
    "internal":     "INTERNAL",
    "confidential": "CONFIDENTIAL",
    "restricted":   "RESTRICTED",
    "secret":       "SECRET",
}

_WATERMARK_COLORS: dict[str, tuple] = {
    "public":       (148, 163, 184),   # Slate — least visible
    "internal":     (251, 191, 36),    # Amber
    "confidential": (239,  68,  68),   # Red
    "restricted":   (220,  38,  38),   # Darker red
    "secret":       (139,   0,   0),   # Dark red — most visible
}


def inject_watermark(
    slide: Slide,
    classification_level: str,
    handling_caveats: list[str] | None = None,
) -> None:
    """
    Inject a diagonal watermark text box across the slide centre.

    The watermark shows the classification label and any handling caveats.
    Text is rotated 315 degrees (= -45 degrees = diagonal from top-right to bottom-left).
    """
    label = _WATERMARK_LABELS.get(classification_level, classification_level.upper())
    if handling_caveats:
        label += "\n" + "  |  ".join(handling_caveats)

    r, g, b = _WATERMARK_COLORS.get(classification_level, (128, 128, 128))

    # Centre the watermark box: 70% of slide width, 30% of slide height
    w = int(_SLIDE_W_EMU * 0.70)
    h = int(_SLIDE_H_EMU * 0.30)
    x = (_SLIDE_W_EMU - w) // 2
    y = (_SLIDE_H_EMU - h) // 2

    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.rotation = 315   # 315° = -45° counterclockwise diagonal

    tf = txBox.text_frame
    tf.word_wrap = False
    tf.clear()

    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label

    run.font.size      = Pt(64)
    run.font.bold      = True
    run.font.color.rgb = RgbColor(r, g, b)
    run.font.name      = "Inter"


def inject_handling_caveats_footer(
    slide: Slide,
    handling_caveats: list[str],
) -> None:
    """
    Add a small italic caveat string to the bottom-left of the slide.
    Used when watermark_required is False but handling_caveats are present.
    """
    if not handling_caveats:
        return

    caveat_text = " | ".join(handling_caveats)

    y = int(_SLIDE_H_EMU * 0.92)
    h = int(_SLIDE_H_EMU * 0.05)
    w = int(_SLIDE_W_EMU * 0.55)
    x = 60_000  # ~60,000 EMU left margin

    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = False
    tf.clear()

    p   = tf.paragraphs[0]
    run = p.add_run()
    run.text       = caveat_text
    run.font.size  = Pt(8)
    run.font.italic = True
    run.font.color.rgb = RgbColor(148, 163, 184)
    run.font.name  = "Inter"
